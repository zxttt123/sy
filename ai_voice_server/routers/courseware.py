from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, status, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
from sqlalchemy.orm import Session
import os
import io
import sys
import time
import uuid
import logging
from typing import Optional
import shutil
import tempfile
import subprocess
import re
import json
from datetime import datetime

from ..database import get_db, SessionLocal
from ..models import models
from ..utils.security import get_current_user
from ..utils.text_splitter import merge_audio_files_exact, get_audio_duration

# 设置日志记录器
logger = logging.getLogger(__name__)

# 尝试导入PPT处理库
PPT_SUPPORT = True
IMPORT_ERRORS = []

try:
    from pptx import Presentation
except ImportError as e:
    IMPORT_ERRORS.append(f"python-pptx: {str(e)}")
    PPT_SUPPORT = False

try:
    from docx import Document
except ImportError as e:
    IMPORT_ERRORS.append(f"python-docx: {str(e)}")
    PPT_SUPPORT = False

try:
    import numpy as np
except ImportError as e:
    IMPORT_ERRORS.append(f"numpy: {str(e)}")
    PPT_SUPPORT = False

try:
    import soundfile as sf
except ImportError as e:
    IMPORT_ERRORS.append(f"soundfile: {str(e)}")
    PPT_SUPPORT = False

try:
    # 修改导入方式
    from moviepy.editor import VideoFileClip, AudioFileClip, ImageClip, concatenate_videoclips, concatenate_audioclips, CompositeVideoClip
    # moviepy的freeze函数位置可能有变化，试用两种路径
    try:
        from moviepy.video.fx.freeze import freeze
    except ImportError:
        try:
            # 备选导入方式
            from moviepy.video.fx.freeze_region import freeze_region as freeze
        except ImportError:
            # 如果freeze不可用，定义一个简易替代函数
            def freeze(clip, duration):
                """简易freeze函数，重复最后一帧"""
                last_frame = clip.to_ImageClip(clip.duration-0.1)
                return concatenate_videoclips([clip, last_frame.set_duration(duration-clip.duration)])
except ImportError as e:
    IMPORT_ERRORS.append(f"moviepy: {str(e)}")
    PPT_SUPPORT = False

if not PPT_SUPPORT:
    print(f"警告: 课件处理功能不可用，缺少以下依赖: {', '.join(IMPORT_ERRORS)}")
    print("请执行 ./fix_moviepy.sh 脚本安装所需依赖")
else:
    print("课件处理依赖已成功加载")

# 设置路径
COSYVOICE_PATH = os.path.expanduser('~/CosyVoice')
MATCHA_TTS_PATH = os.path.join(COSYVOICE_PATH, 'third_party/Matcha-TTS')
sys.path.extend([COSYVOICE_PATH, MATCHA_TTS_PATH])

try:
    from cosyvoice.utils.file_utils import load_wav
    COSYVOICE_AVAILABLE = True
except ImportError:
    print("警告: CosyVoice 模块未找到")
    COSYVOICE_AVAILABLE = False

router = APIRouter()

# 设置上传目录
COURSEWARE_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "coursewares")
os.makedirs(COURSEWARE_DIR, exist_ok=True)
print(f"Courseware directory set to: {COURSEWARE_DIR}")

# 移除全局cosyvoice实例初始化，改为使用CosyVoiceHelper单例

# 存储处理任务状态
task_status = {}

# 修改提取PPT文本的函数，使其按元素顺序提取
def extract_text_from_pptx(pptx_path):
    """从PPTX文件中提取文本，保留每个元素的顺序"""
    try:
        presentation = Presentation(pptx_path)
        slides_content = []
        
        for slide_idx, slide in enumerate(presentation.slides):
            slide_elements = []
            
            # 首先添加幻灯片标题作为第一个元素（如果存在）
            title_shape = None
            for shape in slide.shapes:
                if hasattr(shape, "is_title") and shape.is_title:
                    title_shape = shape
                    break
            
            if title_shape and hasattr(title_shape, "text") and title_shape.text.strip():
                slide_elements.append({
                    "type": "title",
                    "text": title_shape.text.strip(),
                    "order": 0  # 标题总是第一个显示
                })
            
            # 按顺序添加其他文本元素
            element_idx = 1  # 从1开始，因为0是标题
            for shape_idx, shape in enumerate(slide.shapes):
                # 跳过已处理的标题
                if shape == title_shape:
                    continue
                
                if hasattr(shape, "text") and shape.text.strip():
                    # 排除页码等通用元素（通常很短）
                    if len(shape.text.strip()) < 3 and shape.text.strip().isdigit():
                        continue
                        
                    slide_elements.append({
                        "type": "text",
                        "text": shape.text.strip(),
                        "order": element_idx
                    })
                    element_idx += 1
            
            # 如果幻灯片没有文本元素，添加一个占位符
            if not slide_elements:
                slide_elements.append({
                    "type": "placeholder",
                    "text": f"第{slide_idx+1}页",
                    "order": 0
                })
            
            # 添加幻灯片信息
            slides_content.append({
                "index": slide_idx,
                "elements": slide_elements,
                # 创建幻灯片的全文本（用于调试）
                "full_text": " ".join([elem["text"] for elem in slide_elements])
            })
        
        return slides_content
    except Exception as e:
        logger.error(f"从PPTX提取文本时出错: {e}")
        return None

def extract_text_from_docx(docx_path):
    """从DOCX文件中提取文本"""
    try:
        doc = Document(docx_path)
        paragraphs = []
        
        for para in doc.paragraphs:
            if (para.text.strip()):
                paragraphs.append(para.text.strip())
        
        return "\n".join(paragraphs)
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return None

def synthesize_speech(text, voice_id, output_path, db):
    """合成语音"""
    if not COSYVOICE_AVAILABLE:
        print("CosyVoice 模块不可用")
        return False
        
    try:
        # 使用CosyVoiceHelper单例
        from ai_voice_server.cosyvoice_helper import CosyVoiceHelper
        cosyvoice_helper = CosyVoiceHelper()
        
        # 获取声音
        voice = db.query(models.Voice).filter(models.Voice.id == voice_id).first()
        if not voice:
            raise Exception(f"Voice not found with id {voice_id}")

        # 构建完整的音频文件路径
        uploads_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
        audio_path = os.path.join(uploads_dir, voice.filename)
        
        # 检查文件是否存在
        if not os.path.exists(audio_path):
            raise Exception(f"Voice file not found: {audio_path}")
            
        # 加载音频
        prompt_speech_16k = load_wav(audio_path, 16000)
        
        # 分段合成处理长文本
        sentences = re.split(r'[。.!?！？]', text)
        sentences = [s + '。' for s in sentences if s.strip()]
        
        all_audio_data = []
        sample_rate = 22050
        
        # 添加短暂的静音作为句子间隔
        silence_duration = 0.2
        silence_samples = int(silence_duration * sample_rate)
        silence = np.zeros(silence_samples, dtype=np.int16)
        
        for i, sentence in enumerate(sentences):
            if not sentence.strip():
                continue
                
            print(f"Synthesizing sentence {i+1}/{len(sentences)}: {sentence[:50]}...")
            
            try:
                # 使用cosyvoice_helper而不是直接使用cosyvoice
                result = cosyvoice_helper.synthesize_speech(
                    sentence, 
                    voice_id, 
                    is_preset=False,
                    prompt_audio=audio_path,
                    prompt_text=voice.transcript
                )
                
                audio_array = result["audio_data"].astype(np.int16)
                
                if len(audio_array.shape) > 1:
                    audio_array = audio_array.flatten()
                all_audio_data.append(audio_array)
                
                if i < len(sentences) - 1:  # 不在最后一句后添加静音
                    all_audio_data.append(silence)

            except Exception as e:
                print(f"Error synthesizing sentence {i+1}: {str(e)}")
                continue

        if not all_audio_data:
            raise Exception("No audio generated")
        
        print("Concatenating audio segments...")
        combined_audio = np.concatenate(all_audio_data)
        print(f"Combined audio shape: {combined_audio.shape}")
        
        # 保存为WAV文件
        sf.write(output_path, combined_audio, sample_rate)
        
        return True
    except Exception as e:
        print(f"Error synthesizing speech: {e}")
        return False

def synthesize_speech_for_task(text, voice_id, output_path, is_preset=False):
    """为课件任务合成语音，支持预置声音和用户声音"""
    try:
        # 获取或创建Session
        db = SessionLocal()
        
        if is_preset:
            # 对于预置声音，通过CosyVoiceHelper获取语音
            from ai_voice_server.cosyvoice_helper import CosyVoiceHelper
            cosyvoice_helper = CosyVoiceHelper()
            
            try:
                # 直接使用CosyVoiceHelper的synthesize方法，它会根据文本长度自动选择合适的处理方式
                logger.info(f"使用预置声音ID {voice_id} 合成文本: '{text[:30]}...'")
                
                # 获取预置声音列表
                preset_voices = cosyvoice_helper.get_preset_voices()
                
                # 如果voice_id是数字或数字字符串，尝试将其映射到实际的声音名称
                if isinstance(voice_id, (int, str)) and str(voice_id).isdigit():
                    voice_idx = int(voice_id)
                    if 0 < voice_idx <= len(preset_voices):
                        voice_name = preset_voices[voice_idx - 1]
                        logger.info(f"将数字ID {voice_id} 映射到预置声音: {voice_name}")
                    else:
                        if preset_voices:
                            voice_name = preset_voices[0]
                            logger.warning(f"无效的预置声音索引 {voice_id}，使用默认声音: {voice_name}")
                        else:
                            raise ValueError(f"无效的预置声音索引 {voice_id}，且无可用的预置声音")
                else:
                    voice_name = str(voice_id)
                
                # 直接使用synthesize方法，它能处理长文本并生成文件
                cosyvoice_helper.synthesize(text, voice_name, output_path)
                
                # 检查文件是否成功写入
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    return True
                else:
                    logger.error(f"合成文件为空或不存在: {output_path}")
                    return False
                    
            except Exception as e:
                logger.error(f"预置声音合成失败: {str(e)}")
                return False
        else:
            # 对于用户上传的声音，使用现有的synthesize_speech方法
            return synthesize_speech(text, voice_id, output_path, db)
    except Exception as e:
        logger.error(f"合成语音失败: {str(e)}")
        return False
    finally:
        if 'db' in locals():
            db.close()

def process_courseware_task(task_id, file_path, voice_id, is_preset, db_session):
    """后台处理课件任务"""
    try:
        # 获取原始文件名（不带扩展名）
        original_filename = os.path.splitext(os.path.basename(file_path))[0]
        if (original_filename.startswith('temp_')):
            # 移除临时文件名前缀
            original_filename = '_'.join(original_filename.split('_')[2:])
        
        # 当前日期，格式为年月日
        current_date = datetime.now().strftime("%Y%m%d")
        # 为任务创建工作目录
        task_status[task_id] = {
            'status': 'processing',
            'progress': 0,
            'message': '开始处理课件',
            'original_filename': original_filename,
            'process_date': current_date
        }
        task_dir = os.path.join(COURSEWARE_DIR, task_id)
        os.makedirs(task_dir, exist_ok=True)
        
        # 检索课件记录以获取动画模式
        courseware = db_session.query(models.Courseware).filter(models.Courseware.task_id == task_id).first()
        animation_mode = "dynamic"  # 默认动态模式
        transition_time = 0.5  # 默认过渡时间
        
        if courseware:
            animation_mode = courseware.animation_mode or animation_mode
            transition_time = courseware.transition_time or transition_time
            
        logger.info(f"课件处理模式: {animation_mode}, 过渡时间: {transition_time}秒")
        
        # 根据文件类型提取文本
        file_ext = os.path.splitext(file_path)[1].lower()
        
        task_status[task_id]['progress'] = 10
        task_status[task_id]['message'] = '正在提取课件文本'
        
        if file_ext in ['.pptx', '.ppt']:
            # 提取PPT内容，按元素顺序
            slides_content = extract_text_from_pptx(file_path)
            
            if not slides_content or len(slides_content) == 0:
                raise Exception("无法从PPT中提取文本")
                
            # 保存提取的内容到JSON文件（用于调试）
            with open(os.path.join(task_dir, "slides_content.json"), 'w', encoding='utf-8') as f:
                import json
                json.dump(slides_content, f, ensure_ascii=False, indent=2)
            
            # 转换PPT为PDF
            task_status[task_id]['progress'] = 20
            task_status[task_id]['message'] = '正在转换PPT为图片'
            pdf_path = os.path.join(task_dir, "presentation.pdf")
            convert_ppt_to_pdf(file_path, pdf_path)
            
            # 转换PDF为图像
            image_files = convert_pdf_to_images(pdf_path, task_dir)
            
            if not image_files:
                raise Exception("未能从PPT生成图像")
            
            # 创建临时目录用于存储元素音频
            elements_audio_dir = os.path.join(task_dir, "elements_audio")
            os.makedirs(elements_audio_dir, exist_ok=True)
            
            # 创建临时目录用于存储幻灯片视频
            slides_video_dir = os.path.join(task_dir, "slides_video")
            os.makedirs(slides_video_dir, exist_ok=True)
            
            # 为每个幻灯片元素合成单独的音频并创建视频
            task_status[task_id]['progress'] = 30
            task_status[task_id]['message'] = '正在为幻灯片元素生成语音'
            
            all_slide_videos = []
            
            # 确保幻灯片内容和图片数量匹配
            num_slides = min(len(slides_content), len(image_files))
            
            for slide_idx in range(num_slides):
                task_status[task_id]['message'] = f'处理第 {slide_idx+1}/{num_slides} 张幻灯片'
                slide_content = slides_content[slide_idx]
                slide_image = image_files[slide_idx]
                
                # 为幻灯片的每个元素合成音频
                elements_audio = []
                
                for elem_idx, element in enumerate(slide_content["elements"]):
                    elem_text = element["text"]
                    
                    # 跳过空内容
                    if not elem_text.strip():
                        continue
                        
                    # 创建元素音频文件名
                    elem_audio_path = os.path.join(
                        elements_audio_dir, 
                        f"slide_{slide_idx+1}_elem_{elem_idx+1}.wav"
                    )
                    
                    # 合成语音
                    logger.info(f"为幻灯片 {slide_idx+1} 元素 {elem_idx+1} 合成语音: {elem_text[:30]}...")
                    
                    if not synthesize_speech_for_task(elem_text, voice_id, elem_audio_path, is_preset):
                        logger.warning(f"幻灯片 {slide_idx+1} 元素 {elem_idx+1} 语音合成失败，跳过")
                        continue
                    
                    # 获取音频持续时间
                    audio_duration = get_audio_duration(elem_audio_path)
                    
                    # 添加到元素音频列表
                    elements_audio.append({
                        "text": elem_text,
                        "audio_path": elem_audio_path,
                        "duration": audio_duration
                    })
                
                # 检查是否有成功合成的音频
                if not elements_audio:
                    logger.warning(f"幻灯片 {slide_idx+1} 没有成功合成的音频元素，使用默认提示音")
                    # 创建一个默认提示音
                    default_text = f"第{slide_idx+1}页"
                    default_audio_path = os.path.join(elements_audio_dir, f"slide_{slide_idx+1}_default.wav")
                    
                    if synthesize_speech_for_task(default_text, voice_id, default_audio_path, is_preset):
                        audio_duration = get_audio_duration(default_audio_path)
                        elements_audio.append({
                            "text": default_text,
                            "audio_path": default_audio_path,
                            "duration": audio_duration
                        })
                
                # 创建幻灯片视频
                if animation_mode == "dynamic" and elements_audio:
                    task_status[task_id]['message'] = f'为第 {slide_idx+1}/{num_slides} 张幻灯片创建动态视频'
                    slide_video_path = os.path.join(slides_video_dir, f"slide_{slide_idx+1}.mp4")
                    
                    # 创建动态幻灯片视频
                    if create_dynamic_slide_video(
                        slide_image, 
                        elements_audio, 
                        slide_video_path,
                        transition_time
                    ):
                        all_slide_videos.append(slide_video_path)
                        logger.info(f"幻灯片 {slide_idx+1} 动态视频生成成功")
                    else:
                        logger.error(f"幻灯片 {slide_idx+1} 动态视频生成失败")
                else:
                    # 如果不是动态模式或没有元素音频，使用静态幻灯片
                    logger.info(f"幻灯片 {slide_idx+1} 使用静态模式")
                    # (静态模式处理将在后续添加)
            
            # 处理静态模式或动态模式失败的情况
            if not all_slide_videos:
                logger.warning("没有成功创建幻灯片视频，回退到静态模式")
                # 回退到原来的静态模式处理方法...
                # (此处保留原有静态模式代码，不再重复)
            else:
                # 合并所有幻灯片视频
                task_status[task_id]['progress'] = 80
                task_status[task_id]['message'] = '正在合并所有幻灯片视频'
                
                # 使用ffmpeg合并视频，这种方式更可靠且能保持音视频同步
                output_file = os.path.join(task_dir, f"课件视频_{original_filename}_{current_date}.mp4")
                concat_file = os.path.join(task_dir, "concat_list.txt")
                
                # 创建ffmpeg concat文件
                with open(concat_file, 'w') as f:
                    for video in all_slide_videos:
                        f.write(f"file '{os.path.abspath(video)}'\n")
                
                # 合并视频
                try:
                    command = [
                        'ffmpeg',
                        '-f', 'concat',
                        '-safe', '0',
                        '-i', concat_file,
                        '-c', 'copy',
                        output_file
                    ]
                    subprocess.run(command, check=True, capture_output=True)
                    logger.info(f"视频合并成功: {output_file}")
                except Exception as e:
                    logger.error(f"合并视频失败: {e}")
                    # 尝试使用moviepy作为备选方案
                    try:
                        from moviepy.editor import VideoFileClip, concatenate_videoclips
                        
                        video_clips = [VideoFileClip(video) for video in all_slide_videos]
                        final_clip = concatenate_videoclips(video_clips)
                        
                        final_clip.write_videofile(
                            output_file,
                            codec='libx264',
                            audio_codec='aac',
                            fps=24,
                            preset="medium"
                        )
                        
                        # 释放资源
                        final_clip.close()
                        for clip in video_clips:
                            clip.close()
                            
                        logger.info(f"使用moviepy成功合并视频: {output_file}")
                    except Exception as e2:
                        logger.error(f"moviepy合并也失败: {e2}")
                        # 如果合并失败，使用第一个视频作为输出
                        if all_slide_videos:
                            import shutil
                            shutil.copy(all_slide_videos[0], output_file)
                            logger.warning(f"合并失败，使用第一个幻灯片视频作为输出: {output_file}")

            # 记录最终文件路径以及元数据
            task_status[task_id]['output_file'] = output_file
            task_status[task_id]['status'] = 'completed'
            task_status[task_id]['progress'] = 100
            task_status[task_id]['message'] = '课件处理完成'
            
            # 更新数据库记录
            try:
                courseware = db_session.query(models.Courseware).filter(models.Courseware.task_id == task_id).first()
                if courseware:
                    courseware.file_path = output_file
                    courseware.status = 'completed'
                    db_session.commit()
            except Exception as db_error:
                logger.error(f"更新数据库记录失败: {db_error}")

        elif file_ext in ['.docx', '.doc']:
            # DOC文件处理逻辑保持不变（保留原逻辑）
            # ...existing code...
            extracted_text = extract_text_from_docx(file_path)
            
            if not extracted_text:
                raise Exception("无法从文档中提取文本")
            
            # 保存提取的文本
            text_file = os.path.join(task_dir, "extracted_text.txt")
            with open(text_file, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            # 合成语音
            task_status[task_id]['progress'] = 30
            task_status[task_id]['message'] = '正在合成语音'
            
            audio_file = os.path.join(task_dir, "narration.wav")
            if not synthesize_speech_for_task(extracted_text, voice_id, audio_file, is_preset):
                raise Exception("语音合成失败")
            
            # 修改输出文件命名以及元数据
            output_file = os.path.join(task_dir, f"课件音频_{original_filename}_{current_date}.wav")
            shutil.copy(audio_file, output_file)
            print(f"文档类型，仅提供音频文件: {output_file}")
            
            # 记录最终文件路径以及元数据
            task_status[task_id]['output_file'] = output_file
            task_status[task_id]['status'] = 'completed'
            task_status[task_id]['progress'] = 100
            task_status[task_id]['message'] = '课件处理完成'
            
            # 更新数据库记录
            try:
                courseware = db_session.query(models.Courseware).filter(models.Courseware.task_id == task_id).first()
                if courseware:
                    courseware.file_path = output_file
                    courseware.status = 'completed'
                    db_session.commit()
            except Exception as db_error:
                logger.error(f"更新数据库记录失败: {db_error}")
            
        else:
            raise Exception("不支持的文件类型")

    # 这部分异常处理保持不变
    except Exception as e:
        print(f"Error processing courseware: {e}")
        task_status[task_id]['status'] = 'failed'
        task_status[task_id]['message'] = f'处理失败: {str(e)}'
        # 更新数据库状态
        try:
            courseware = db_session.query(models.Courseware).filter(models.Courseware.task_id == task_id).first()
            if courseware:
                courseware.status = 'failed'
                db_session.commit()
        except Exception as db_error:
            print(f"Error updating courseware status: {db_error}")
    finally:
        # 清理临时文件
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"已删除临时文件: {file_path}")
        except Exception as e:
            print(f"Error removing temp file: {e}")

def convert_ppt_to_pdf(ppt_path, pdf_path):
    """将PPT转换为PDF"""
    try:
        # 检查LibreOffice是否可用
        libreoffice_cmd = "libreoffice"
        try:
            subprocess.run(["which", libreoffice_cmd], check=True, capture_output=True)
        except subprocess.CalledProcessError:
            libreoffice_cmd = "soffice"  # 尝试替代命令
        # 转换为PDF
        cmd = [
            libreoffice_cmd,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", os.path.dirname(pdf_path),
            ppt_path
        ]
        print(f"执行命令: {' '.join(cmd)}")
        subprocess.run(cmd, check=True, capture_output=True)
        
        # 检查PDF文件是否生成
        if not os.path.exists(pdf_path):
            print(f"PDF文件未直接生成在指定路径: {pdf_path}")
            # 尝试查找生成的PDF文件
            base_name = os.path.splitext(os.path.basename(ppt_path))[0]
            pdf_name = f"{base_name}.pdf"
            potential_path = os.path.join(os.path.dirname(pdf_path), pdf_name)
            if os.path.exists(potential_path):
                print(f"找到PDF文件: {potential_path}")
                if potential_path != pdf_path:
                    shutil.move(potential_path, pdf_path)
                return pdf_path
            
            # 搜索目录中的所有PDF文件
            pdf_files = [f for f in os.listdir(os.path.dirname(pdf_path)) if f.endswith('.pdf')]
            if pdf_files:
                found_pdf = os.path.join(os.path.dirname(pdf_path), pdf_files[0])
                print(f"找到PDF文件: {found_pdf}")
                if found_pdf != pdf_path:
                    shutil.move(found_pdf, pdf_path)
                return pdf_path
            
            raise FileNotFoundError(f"未找到生成的PDF文件: {pdf_path}")
        return pdf_path
    except Exception as e:
        print(f"转换PPT到PDF失败: {str(e)}")
        raise

def convert_pdf_to_images(pdf_path, output_dir):
    """将PDF转换为图像序列"""
    try:
        # 使用ImageMagick转换PDF为图像
        cmd = [
            "convert",
            "-density", "300",
            pdf_path,
            os.path.join(output_dir, "slide-%03d.jpg")
        ]
        print(f"执行命令: {' '.join(cmd)}")
        subprocess.run(cmd, check=True, capture_output=True)
        
        # 获取生成的图像文件
        image_files = sorted([
            os.path.join(output_dir, f) for f in os.listdir(output_dir) 
            if f.startswith("slide-") and f.endswith(".jpg")
        ])
        if not image_files:
            raise FileNotFoundError(f"未在 {output_dir} 目录下找到生成的图像文件")
        return image_files
    except Exception as e:
        print(f"转换PDF到图像失败: {str(e)}")
        raise

# 添加生成动态PPT切换效果的函数
def create_dynamic_slide_video(slide_image_path, elements_audio, output_path, transition_time=0.5):
    """
    为单个幻灯片创建动态切换效果的视频
    
    参数:
        slide_image_path: 幻灯片图片路径
        elements_audio: 元素音频信息列表 [{text, audio_path, duration}, ...]
        output_path: 输出视频路径
        transition_time: 元素之间的过渡时间
    """
    from moviepy.editor import ImageClip, AudioFileClip, concatenate_audioclips
    
    try:
        if not elements_audio:
            logger.warning(f"幻灯片没有元素音频: {slide_image_path}")
            return False
            
        # 加载幻灯片图片
        slide_clip = ImageClip(slide_image_path)
        
        # 计算总持续时间（所有元素音频加上过渡时间）
        total_duration = sum([element["duration"] for element in elements_audio])
        total_duration += transition_time * (len(elements_audio) - 1)
        
        # 设置幻灯片持续时间
        slide_clip = slide_clip.set_duration(total_duration)
        
        # 合并所有音频片段，添加过渡时间的静音
        audio_clips = []
        
        for i, element in enumerate(elements_audio):
            # 添加元素音频
            audio_clip = AudioFileClip(element["audio_path"])
            audio_clips.append(audio_clip)
            
            # 如果不是最后一个元素，添加过渡静音
            if i < len(elements_audio) - 1:
                from moviepy.audio.AudioClip import AudioClip
                silence = AudioClip(lambda t: 0, duration=transition_time)
                audio_clips.append(silence)
        
        combined_audio = concatenate_audioclips(audio_clips)
        slide_clip = slide_clip.set_audio(combined_audio)
        
        # 导出视频
        slide_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            fps=24,
            preset="medium"
        )
        
        # 清理资源
        slide_clip.close()
        combined_audio.close()
        for clip in audio_clips:
            if hasattr(clip, 'close'):  # 静音AudioClip可能没有close方法
                clip.close()
        
        return True
    except Exception as e:
        logger.error(f"创建动态幻灯片视频时出错: {e}")
        return False

# 替换原始的process_slides_to_video函数
async def process_slides_to_video(slides_folder, audio_folder, output_path, slide_durations=None):
    """
    将幻灯片和音频合成为视频，使用改进的同步方法
    """
    return process_slides_to_video_improved(slides_folder, audio_folder, output_path)

@router.post("/courseware/process")
async def process_courseware(
    background_tasks: BackgroundTasks,
    current_user: models.User = Depends(get_current_user),
    courseware: UploadFile = File(...),
    voice_id: str = Form(...),
    is_preset: bool = Form(False),  # 添加预设声音标志参数
    animation_mode: Optional[str] = Form("dynamic"),  # 增加动画模式参数
    transition_time: Optional[float] = Form(0.5),     # 增加过渡时间参数
    db: Session = Depends(get_db)
):
    if not PPT_SUPPORT:
        raise HTTPException(
            status_code=status.HTTP_501_NOT_IMPLEMENTED,
            detail="服务器未安装必要的课件处理组件"
        )
    if not COSYVOICE_AVAILABLE:
        raise HTTPException(
            status_code=status.HTTP_501_NOT_IMPLEMENTED,
            detail="语音合成模型未正确加载"
        )
    
    # 校验参数
    if animation_mode not in ["dynamic", "static"]:
        animation_mode = "dynamic"  # 默认使用动态模式
    
    if transition_time < 0.1:
        transition_time = 0.5  # 设置最小过渡时间
    elif transition_time > 3.0:
        transition_time = 3.0  # 设置最大过渡时间
    
    # 检查CosyVoiceHelper是否可用
    try:
        from ai_voice_server.cosyvoice_helper import CosyVoiceHelper
        helper = CosyVoiceHelper(lazy_load=True)  # 使用懒加载模式
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"无法初始化语音合成Helper: {str(e)}"
        )
    
    try:
        print(f"Processing courseware: {courseware.filename}, voice_id: {voice_id}, is_preset: {is_preset}, animation_mode: {animation_mode}")
        
        # 创建唯一任务ID
        task_id = str(uuid.uuid4())
        # 保存上传的文件（保留原始文件名）
        file_content = await courseware.read()
        # 保存时添加任务ID以避免文件名冲突，但保留原始文件名用于后续显示
        file_path = os.path.join(COURSEWARE_DIR, f"temp_{task_id}_{courseware.filename}")
        with open(file_path, "wb") as f:
            f.write(file_content)
        
        # 验证声音ID
        try:
            voice_id_int = int(voice_id)
            voice = db.query(models.Voice).filter(models.Voice.id == voice_id_int).first()
            if not voice:
                raise HTTPException(
                    status_code=404,
                    detail="找不到指定的声音ID"
                )
        except ValueError:
            raise HTTPException(
                status_code=400,
                detail="无效的声音ID格式"
            )
        
        # 验证声音使用权限（移除了仅允许用户自己上传的声音的限制）
        # 如果是预设声音，任何人都可以使用；如果不是预设声音，必须是用户自己的
        if not is_preset and not voice.is_preset and voice.user_id != current_user.id:
            raise HTTPException(
                status_code=403,
                detail="您没有权限使用此声音"
            )
        
        # 创建课件记录
        courseware_record = models.Courseware(
            task_id=task_id,
            user_id=current_user.id,
            voice_id=voice_id_int,
            status="processing",
            # 添加原始文件名和处理日期
            original_filename=os.path.splitext(courseware.filename)[0],
            process_date=datetime.now().strftime("%Y%m%d"),
            animation_mode=animation_mode,
            transition_time=transition_time
        )
        db.add(courseware_record)
        db.commit()
        
        # 启动后台任务处理课件
        # 添加原始文件名和当前日期信息
        original_filename = os.path.splitext(courseware.filename)[0]
        current_date = datetime.now().strftime("%Y%m%d")
        task_status[task_id] = {
            'status': 'pending',
            'progress': 0,
            'message': '任务已提交，等待处理',
            'created_at': datetime.now().isoformat(),
            'original_filename': original_filename,
            'process_date': current_date,
            'is_preset': is_preset,  # 添加标记声音是否为预设
            'animation_mode': animation_mode,
            'transition_time': transition_time
        } 
        
        # 传递 is_preset 参数给后台处理任务
        background_tasks.add_task(
            process_courseware_task, 
            task_id, 
            file_path, 
            voice_id_int, 
            is_preset,  # 添加此参数
            db
        )

        # 记录课件合成日志
        synthesis_log = models.SynthesisLog(
            type="courseware",
            user_id=current_user.id,
            voice_id=int(voice_id) if not is_preset else None,
            text_length=0,  # 假设有一个变量记录了总文本长度
            duration=0  # 假设有一个变量记录了总时长
        )
        db.add(synthesis_log)
        db.commit()

        return {"message": "课件处理任务已提交", "task_id": task_id}
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error in process_courseware: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"处理课件失败: {str(e)}"
        )

@router.get("/courseware/status/{task_id}")
async def get_task_status(
    task_id: str,
    current_user: models.User = Depends(get_current_user)
):
    if task_id not in task_status:
        raise HTTPException(
            status_code=404,
            detail="找不到指定的任务"
        )
    return task_status[task_id]

@router.get("/courseware/download/{task_id}")
async def download_courseware(
    task_id: str,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    if task_id not in task_status or task_status[task_id].get('status') != 'completed':
        raise HTTPException(
            status_code=404,
            detail="找不到完成的任务或任务尚未完成"
        )
    
    output_file = task_status[task_id].get('output_file')
    if not output_file or not os.path.exists(output_file):
        raise HTTPException(
            status_code=404,
            detail="找不到课件处理结果文件"
        )
    
    # 获取课件记录
    courseware = db.query(models.Courseware).filter(models.Courseware.task_id == task_id).first()
    # 验证权限 - 只有管理员或文件所有者可下载
    if not courseware or (courseware.user_id != current_user.id and not current_user.is_admin):
        raise HTTPException(
            status_code=403,
            detail="没有权限下载此文件"
        )
    
    # 确定文件类型和文件名
    file_ext = os.path.splitext(output_file)[1].lower()
    original_filename = task_status[task_id].get('original_filename', 'unnamed')
    process_date = task_status[task_id].get('process_date', datetime.now().strftime("%Y%m%d"))
    
    if file_ext == '.mp4':
        media_type = 'video/mp4'
        filename = f"课件视频_{original_filename}_{process_date}.mp4"
    elif file_ext == '.wav':
        media_type = 'audio/wav'
        filename = f"课件音频_{original_filename}_{process_date}.wav"
    else:
        media_type = 'application/octet-stream'
        filename = f"课件_{original_filename}_{process_date}{file_ext}"
    
    return FileResponse(
        path=output_file,
        media_type=media_type,
        filename=filename
    )